"""오피스 자동화 도구 - PPT·Excel·한글(DOCX) 파일을 자동 생성한다.

LLM이 내용을 구조화해 전달하면 이 모듈이 실제 파일을 생성한다.
생성된 파일은 지정 경로에 저장되고, 기본값은 사용자 Documents 폴더다.
"""

import os
from pathlib import Path
from datetime import datetime

from config import settings

# ── 기본 저장 경로 ─────────────────────────────────────────────────────────────
DEFAULT_OUTPUT_DIR = settings.DOCUMENTS_DIR / "AgentAI_문서"


def _resolve_output_path(output_path: str, default_name: str) -> Path:
    """출력 경로를 결정한다. 미지정 시 Documents/AgentAI_문서 폴더에 저장한다."""
    if output_path:
        p = Path(output_path)
    else:
        DEFAULT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        p = DEFAULT_OUTPUT_DIR / default_name
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _open_file(path: Path) -> None:
    """생성된 파일을 기본 프로그램으로 바로 연다 (Windows)."""
    try:
        os.startfile(str(path))
    except Exception:
        pass  # 파일 열기 실패는 무시 (생성 자체는 성공)


# ══════════════════════════════════════════════════════════════════════════════
# PPT 생성
# ══════════════════════════════════════════════════════════════════════════════

def create_pptx(
    title: str,
    slides: list,
    output_path: str = "",
    theme: str = "default",
    open_after: bool = True,
) -> dict:
    """PowerPoint 파일을 자동 생성한다.

    Args:
        title:       프레젠테이션 전체 제목
        slides:      슬라이드 목록. 각 항목은 dict:
                       {"제목": "...", "내용": ["항목1", "항목2", ...],
                        "레이아웃": "title|content|two_column|blank",
                        "노트": "발표자 노트 (선택)"}
        output_path: 저장 경로 (.pptx). 비워두면 Documents/AgentAI_문서에 저장.
        theme:       "default"(파랑) | "dark"(다크) | "minimal"(미니멀)
        open_after:  생성 후 파일 바로 열기
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # ── 테마 색상 ────────────────────────────────────────────────────────
        THEMES = {
            "default": {
                "title_bg":   RGBColor(0x1F, 0x49, 0x7D),
                "content_bg": RGBColor(0xFF, 0xFF, 0xFF),
                "title_fg":   RGBColor(0xFF, 0xFF, 0xFF),
                "content_fg": RGBColor(0x1F, 0x1F, 0x1F),
                "accent":     RGBColor(0x2E, 0x75, 0xB6),
                "bullet_fg":  RGBColor(0x2E, 0x75, 0xB6),
            },
            "dark": {
                "title_bg":   RGBColor(0x1A, 0x1A, 0x2E),
                "content_bg": RGBColor(0x16, 0x21, 0x3E),
                "title_fg":   RGBColor(0xE9, 0x4C, 0x7C),
                "content_fg": RGBColor(0xE0, 0xE0, 0xE0),
                "accent":     RGBColor(0xE9, 0x4C, 0x7C),
                "bullet_fg":  RGBColor(0xE9, 0x4C, 0x7C),
            },
            "minimal": {
                "title_bg":   RGBColor(0xF5, 0xF5, 0xF5),
                "content_bg": RGBColor(0xFF, 0xFF, 0xFF),
                "title_fg":   RGBColor(0x22, 0x22, 0x22),
                "content_fg": RGBColor(0x44, 0x44, 0x44),
                "accent":     RGBColor(0x00, 0x78, 0xD7),
                "bullet_fg":  RGBColor(0x00, 0x78, 0xD7),
            },
        }
        c = THEMES.get(theme, THEMES["default"])

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 슬레이드 레이아웃 인덱스: 0=표지, 1=제목+내용, 6=빈 슬라이드
        blank_layout   = prs.slide_layouts[6]

        def _set_bg(slide, color: RGBColor):
            from pptx.oxml.ns import qn
            from lxml import etree
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_textbox(slide, text, left, top, width, height,
                         font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
            from pptx.util import Emu
            txb = slide.shapes.add_textbox(left, top, width, height)
            txb.text_frame.word_wrap = wrap
            p = txb.text_frame.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            return txb

        # ── 표지 슬라이드 (첫 번째) ──────────────────────────────────────────
        cover = prs.slides.add_slide(blank_layout)
        _set_bg(cover, c["title_bg"])
        # 장식 사각형 (하단 강조선)
        bar = cover.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(6.8), Inches(13.33), Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = c["accent"]
        bar.line.fill.background()

        _add_textbox(cover, title,
                     Inches(1), Inches(2.5), Inches(11), Inches(1.8),
                     font_size=40, bold=True, color=c["title_fg"],
                     align=PP_ALIGN.CENTER)
        timestamp = datetime.now().strftime("%Y.%m.%d")
        _add_textbox(cover, timestamp,
                     Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                     font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
                     align=PP_ALIGN.CENTER)

        # ── 본문 슬라이드 ────────────────────────────────────────────────────
        for slide_data in slides:
            layout_type = slide_data.get("레이아웃", "content")
            slide_title  = slide_data.get("제목", "")
            contents     = slide_data.get("내용", [])
            notes_text   = slide_data.get("노트", "")

            slide = prs.slides.add_slide(blank_layout)
            _set_bg(slide, c["content_bg"])

            # 상단 제목 영역 배경
            title_bar = slide.shapes.add_shape(
                1,
                Inches(0), Inches(0), Inches(13.33), Inches(1.2)
            )
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = c["accent"]
            title_bar.line.fill.background()

            # 슬라이드 제목
            _add_textbox(slide, slide_title,
                         Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.9),
                         font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

            if layout_type == "title":
                # 큰 텍스트 한 줄
                body = contents[0] if contents else ""
                _add_textbox(slide, body,
                             Inches(1), Inches(2), Inches(11), Inches(3),
                             font_size=32, bold=True, color=c["content_fg"],
                             align=PP_ALIGN.CENTER)

            elif layout_type == "two_column":
                # 두 칼럼: 홀수 항목 왼쪽 / 짝수 항목 오른쪽
                left_items  = contents[::2]
                right_items = contents[1::2]
                for items, left_offset in [(left_items, 0.4), (right_items, 6.9)]:
                    txb = slide.shapes.add_textbox(
                        Inches(left_offset), Inches(1.4), Inches(6.0), Inches(5.5)
                    )
                    txb.text_frame.word_wrap = True
                    for i, item in enumerate(items):
                        p = txb.text_frame.add_paragraph() if i > 0 else txb.text_frame.paragraphs[0]
                        p.space_before = Pt(6)
                        run = p.add_run()
                        run.text = f"• {item}"
                        run.font.size = Pt(18)
                        run.font.color.rgb = c["content_fg"]

            else:  # "content" — 기본 불릿 리스트
                txb = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6)
                )
                txb.text_frame.word_wrap = True
                for i, item in enumerate(contents):
                    p = txb.text_frame.add_paragraph() if i > 0 else txb.text_frame.paragraphs[0]
                    p.space_before = Pt(8)
                    run = p.add_run()
                    # 들여쓰기(하위 항목)는 "  " 두 칸 접두사로 표현
                    if str(item).startswith("  "):
                        run.text = f"    ◦ {item.strip()}"
                        run.font.size = Pt(17)
                        run.font.color.rgb = c["content_fg"]
                    else:
                        run.text = f"  ▸ {item}"
                        run.font.size = Pt(20)
                        run.font.bold = False
                        run.font.color.rgb = c["bullet_fg"]

            # 발표자 노트
            if notes_text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text

        # ── 저장 ─────────────────────────────────────────────────────────────
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        default_name = f"{ts}_{title[:20]}.pptx"
        path = _resolve_output_path(output_path, default_name)
        if path.suffix.lower() != ".pptx":
            path = path.with_suffix(".pptx")

        prs.save(str(path))
        if open_after:
            _open_file(path)

        return {
            "결과": "성공",
            "파일_경로": str(path),
            "슬라이드_수": len(slides) + 1,  # 표지 포함
        }
    except ImportError:
        return {"결과": "실패", "이유": "python-pptx 가 설치되어 있지 않습니다. pip install python-pptx"}
    except Exception as e:
        return {"결과": "실패", "이유": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# Excel 생성
# ══════════════════════════════════════════════════════════════════════════════

def create_excel(
    title: str,
    sheets: list,
    output_path: str = "",
    open_after: bool = True,
) -> dict:
    """Excel 파일을 자동 생성한다.

    Args:
        title:       파일 대표 제목 (첫 번째 시트 상단에 표시)
        sheets:      시트 목록. 각 항목은 dict:
                       {"시트명": "...",
                        "헤더": ["열1", "열2", ...],
                        "데이터": [["값1", "값2"], ["값3", "값4"], ...],
                        "수식": {"B10": "=SUM(B2:B9)", ...},   (선택)
                        "너비": {"A": 20, "B": 15, ...}}       (선택)
        output_path: 저장 경로 (.xlsx)
        open_after:  생성 후 바로 열기
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import (Font, PatternFill, Alignment,
                                     Border, Side, GradientFill)
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        wb.remove(wb.active)  # 기본 시트 제거

        # 공통 스타일
        HEADER_FILL = PatternFill("solid", fgColor="1F497D")
        HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
        TITLE_FONT  = Font(bold=True, size=14, color="1F497D")
        ALT_FILL    = PatternFill("solid", fgColor="EBF3FB")
        BORDER_SIDE = Side(style="thin", color="CCCCCC")
        CELL_BORDER = Border(
            left=BORDER_SIDE, right=BORDER_SIDE,
            top=BORDER_SIDE,  bottom=BORDER_SIDE
        )
        CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
        LEFT   = Alignment(horizontal="left",   vertical="center", wrap_text=True)

        for idx, sheet_data in enumerate(sheets):
            sheet_name = sheet_data.get("시트명", f"시트{idx + 1}")
            headers    = sheet_data.get("헤더", [])
            data_rows  = sheet_data.get("데이터", [])
            formulas   = sheet_data.get("수식", {})
            col_widths = sheet_data.get("너비", {})

            ws = wb.create_sheet(title=sheet_name)

            row_offset = 1

            # 제목 행 (첫 번째 시트만)
            if idx == 0 and title:
                ws.merge_cells(
                    start_row=1, start_column=1,
                    end_row=1,   end_column=max(len(headers), 1)
                )
                ws.cell(1, 1).value = title
                ws.cell(1, 1).font = TITLE_FONT
                ws.cell(1, 1).alignment = CENTER
                ws.row_dimensions[1].height = 28
                row_offset = 2

            # 헤더 행
            if headers:
                for col, h in enumerate(headers, 1):
                    cell = ws.cell(row_offset, col)
                    cell.value = h
                    cell.font = HEADER_FONT
                    cell.fill = HEADER_FILL
                    cell.alignment = CENTER
                    cell.border = CELL_BORDER
                ws.row_dimensions[row_offset].height = 22
                row_offset += 1

            # 데이터 행
            for r_idx, row in enumerate(data_rows):
                fill = ALT_FILL if r_idx % 2 == 1 else None
                for col, val in enumerate(row, 1):
                    cell = ws.cell(row_offset + r_idx, col)
                    cell.value = val
                    cell.alignment = LEFT
                    cell.border = CELL_BORDER
                    if fill:
                        cell.fill = fill
                ws.row_dimensions[row_offset + r_idx].height = 18

            # 수식 추가
            for cell_addr, formula in formulas.items():
                ws[cell_addr] = formula
                ws[cell_addr].font = Font(bold=True, color="1F497D")
                ws[cell_addr].border = CELL_BORDER

            # 열 너비 조정
            if col_widths:
                for col_letter, width in col_widths.items():
                    ws.column_dimensions[col_letter.upper()].width = width
            elif headers:
                # 자동 너비 추정
                for col, h in enumerate(headers, 1):
                    letter = get_column_letter(col)
                    max_len = len(str(h))
                    for row in data_rows:
                        if col - 1 < len(row):
                            max_len = max(max_len, len(str(row[col - 1])))
                    ws.column_dimensions[letter].width = min(max_len + 4, 40)

            # 틀 고정 (헤더 행 아래)
            if headers:
                ws.freeze_panes = ws.cell(row_offset, 1)

            # 자동 필터
            if headers:
                ws.auto_filter.ref = ws.dimensions

        # ── 저장 ─────────────────────────────────────────────────────────────
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        default_name = f"{ts}_{title[:20]}.xlsx"
        path = _resolve_output_path(output_path, default_name)
        if path.suffix.lower() != ".xlsx":
            path = path.with_suffix(".xlsx")

        wb.save(str(path))
        if open_after:
            _open_file(path)

        total_rows = sum(len(s.get("데이터", [])) for s in sheets)
        return {
            "결과": "성공",
            "파일_경로": str(path),
            "시트_수": len(sheets),
            "총_데이터_행": total_rows,
        }
    except Exception as e:
        return {"결과": "실패", "이유": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# 한글·Word 문서 생성 (DOCX)
# ══════════════════════════════════════════════════════════════════════════════

def create_document(
    title: str,
    sections: list,
    output_path: str = "",
    open_after: bool = True,
) -> dict:
    """한글/Word 문서(.docx)를 자동 생성한다. HWP에서도 열 수 있다.

    Args:
        title:        문서 제목
        sections:     섹션 목록. 각 항목은 dict:
                        {"제목": "소제목",    (선택 — 없으면 본문만)
                         "레벨": 1,          (1=대제목 / 2=소제목 / 3=소소제목)
                         "내용": "본문 텍스트...",
                         "표": {             (선택 — 표 삽입)
                           "헤더": ["열1","열2"],
                           "데이터": [["a","b"],["c","d"]]
                         },
                         "목록": ["항목1","항목2"]}  (선택 — 불릿 리스트)
        output_path:  저장 경로 (.docx)
        open_after:   생성 후 바로 열기
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()

        # ── 기본 스타일 설정 ─────────────────────────────────────────────────
        style = doc.styles["Normal"]
        style.font.name = "맑은 고딕"
        style.font.size = Pt(10.5)

        # 페이지 여백
        for section in doc.sections:
            section.top_margin    = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin   = Cm(3.0)
            section.right_margin  = Cm(3.0)

        # ── 문서 제목 ────────────────────────────────────────────────────────
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_para.runs[0] if title_para.runs else title_para.add_run(title)
        title_run.font.size = Pt(18)
        title_run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        title_run.font.bold = True
        title_run.font.name = "맑은 고딕"

        doc.add_paragraph()  # 빈 줄

        # ── 섹션 추가 ────────────────────────────────────────────────────────
        for sec in sections:
            heading_text = sec.get("제목", "")
            level        = sec.get("레벨", 1)
            body_text    = sec.get("내용", "")
            table_data   = sec.get("표", None)
            bullet_items = sec.get("목록", [])

            # 소제목
            if heading_text:
                h = doc.add_heading(heading_text, level=level)
                if h.runs:
                    run = h.runs[0]
                    run.font.name = "맑은 고딕"
                    if level == 1:
                        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
                        run.font.size = Pt(14)
                    elif level == 2:
                        run.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
                        run.font.size = Pt(12)
                    else:
                        run.font.size = Pt(11)

            # 본문 텍스트
            if body_text:
                for line in body_text.split("\n"):
                    p = doc.add_paragraph(line)
                    p.paragraph_format.space_after  = Pt(4)
                    p.paragraph_format.line_spacing = Pt(18)
                    if p.runs:
                        p.runs[0].font.name = "맑은 고딕"

            # 불릿 리스트
            for item in bullet_items:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(item)
                run.font.name = "맑은 고딕"
                run.font.size = Pt(10.5)

            # 표
            if table_data:
                t_headers = table_data.get("헤더", [])
                t_rows    = table_data.get("데이터", [])
                if t_headers:
                    table = doc.add_table(
                        rows=1 + len(t_rows),
                        cols=len(t_headers)
                    )
                    table.style = "Table Grid"

                    # 헤더 행
                    hdr_cells = table.rows[0].cells
                    for i, h_text in enumerate(t_headers):
                        hdr_cells[i].text = h_text
                        run = hdr_cells[i].paragraphs[0].runs[0]
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.name = "맑은 고딕"
                        # 헤더 배경색
                        tc = hdr_cells[i]._tc
                        tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement("w:shd")
                        shd.set(qn("w:fill"), "1F497D")
                        shd.set(qn("w:color"), "auto")
                        shd.set(qn("w:val"),  "clear")
                        tcPr.append(shd)

                    # 데이터 행
                    for r_idx, row_vals in enumerate(t_rows):
                        row_cells = table.rows[r_idx + 1].cells
                        for c_idx, val in enumerate(row_vals):
                            row_cells[c_idx].text = str(val)
                            if row_cells[c_idx].paragraphs[0].runs:
                                row_cells[c_idx].paragraphs[0].runs[0].font.name = "맑은 고딕"

                doc.add_paragraph()  # 표 뒤 빈 줄

        # ── 저장 ─────────────────────────────────────────────────────────────
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        default_name = f"{ts}_{title[:20]}.docx"
        path = _resolve_output_path(output_path, default_name)
        if path.suffix.lower() != ".docx":
            path = path.with_suffix(".docx")

        doc.save(str(path))
        if open_after:
            _open_file(path)

        return {
            "결과": "성공",
            "파일_경로": str(path),
            "섹션_수": len(sections),
            "안내": "한글(HWP) 2020 이상에서도 .docx 파일을 바로 열 수 있습니다.",
        }
    except Exception as e:
        return {"결과": "실패", "이유": str(e)}
